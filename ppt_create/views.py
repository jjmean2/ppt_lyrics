import os
import re
from django.shortcuts import render
from pptx import Presentation
from io import BytesIO
from django.http import HttpResponse
from .forms import LyricsForm
from django.views.decorators.csrf import csrf_exempt
from urllib.parse import quote


@csrf_exempt
def lyrics(request):
    form = LyricsForm()
    if request.method == 'POST':
        form = LyricsForm(request.POST)
        if form.is_valid():
            body = form.cleaned_data['body']
            filename = form.cleaned_data['filename']
            ppt = create_ppt_content(body)
            filename_header = 'filename*=UTF-8\'\'%s' % quote(
                filename.encode('utf-8'))
            response = HttpResponse(
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = 'attachment; ' + filename_header
            response.write(ppt)
            return response
    else:
        context = {
            'form': form
        }
        return render(request, 'lyrics.html', context)


def create_ppt_content(body):
    module_dir = os.path.dirname(__file__)
    template_file = 'resources/tehilla_ppt_template2.pptx'
    file_path = os.path.join(module_dir, template_file)
    prs = Presentation(file_path)
    fill_lyrics_slides(prs, body)
    source_stream = BytesIO()
    prs.save(source_stream)
    ppt = source_stream.getvalue()
    source_stream.close()
    return ppt


def get_slide_sources_from_text(body):
    sanitizedBody = '\n'.join(body.splitlines())
    print(sanitizedBody)
    songs = re.split(r'\n={3,}\n', sanitizedBody)
    sources = []
    for song in songs:
        pageSources = []
        pages = re.split(r'\n-{3,}\n', song)
        for page in pages:
            lines = page.splitlines()
            tag = lines[0].lstrip('[').rstrip(']')
            content = '\n'.join(lines[1:])
            pageSources.append((tag, content))
        sources.append(pageSources)
    print('sources', sources)
    return sources


def fill_lyrics_slides(prs, body):
    lyrics_slide_layout = prs.slide_layouts[1]
    empty_slide_layout = prs.slide_layouts[2]
    TITLE_PLACEHOLDER_IDX = 0
    TAG_PLACEHOLDER_IDX = 10

    sources = get_slide_sources_from_text(body)

    for pageSources in sources:
        for (tag, content) in pageSources:
            slide = prs.slides.add_slide(lyrics_slide_layout)
            title_box = slide.placeholders[TITLE_PLACEHOLDER_IDX]
            title_box.text = content
            tag_box = slide.placeholders[TAG_PLACEHOLDER_IDX]
            tag_box.text = tag
        prs.slides.add_slide(empty_slide_layout)
